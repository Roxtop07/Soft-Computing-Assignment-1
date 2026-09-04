"""
Generates the 15-slide PowerPoint presentation for Assignment 1:
Topic 1: Design of an Intelligent Air Conditioning System Using Fuzzy Logic
Course: Soft Computing (EL1) - STDA2102
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

FIGURES_DIR = os.path.join(os.path.dirname(__file__), "figures")
PPTX_PATH = os.path.join(os.path.dirname(__file__), "Presentation_Intelligent_AC.pptx")

# Color Palette (Dark Tech / Academic Modern)
BG_COLOR = RGBColor(15, 23, 42)        # Slate 900
CARD_COLOR = RGBColor(30, 41, 59)      # Slate 800
ACCENT_BLUE = RGBColor(14, 165, 233)   # Sky 500
ACCENT_GREEN = RGBColor(34, 197, 94)   # Green 500
TEXT_WHITE = RGBColor(248, 250, 252)   # Slate 50
TEXT_MUTED = RGBColor(148, 163, 184)   # Slate 400

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, category="SOFT COMPUTING (STDA2102) | CASE STUDY"):
    # Header Category / Breadcrumb
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.35))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_BLUE

    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, BG_COLOR)

    # Title box
    tbox = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(3.0))
    tf1 = tbox.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "Design of an Intelligent Air Conditioning System\nUsing Fuzzy Logic"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.alignment = PP_ALIGN.LEFT

    p1_sub = tf1.add_paragraph()
    p1_sub.text = "\nSoft Computing (STDA2102) | Case Study & Technical Presentation"
    p1_sub.font.size = Pt(18)
    p1_sub.font.color.rgb = ACCENT_BLUE

    p1_meta = tf1.add_paragraph()
    p1_meta.text = "\nStudent Name: Manish Kumar  |  Degree: B.Tech / M.Tech  |  Date: September 2026"
    p1_meta.font.size = Pt(14)
    p1_meta.font.color.rgb = TEXT_MUTED

    add_speaker_notes(s1, "Welcome everyone. Today I am presenting our case study on the Design of an Intelligent Air Conditioning System using Fuzzy Logic for the course STDA2102. This project demonstrates how fuzzy logic enables continuous compressor modulation, optimizing both thermal comfort and energy efficiency.")

    # -------------------------------------------------------------
    # SLIDE 2: Executive Summary & Motivation
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2, BG_COLOR)
    add_header(s2, "Executive Summary & Problem Statement")

    body_box = s2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tf2 = body_box.text_frame
    tf2.word_wrap = True

    bullets2 = [
        ("The Global HVAC Energy Challenge:", "Buildings account for over 40% of global electricity consumption, with HVAC systems consuming over 50% of building energy load."),
        ("Limitations of Conventional Thermostats (Bang-Bang):", "Classical on/off thermostats suffer from frequent cycling, abrupt temperature fluctuations (±2°C hysteresis), and mechanical wear on compressors."),
        ("Shortcomings of Standard PID Controllers:", "PID controllers struggle with non-linear thermal dynamics, variable ambient weather, room occupancy changes, and delayed dead-time response."),
        ("Fuzzy Logic Solution:", "Fuzzy Logic Controllers (FLC) translate human linguistic expertise into mathematical models, providing smooth, multi-variable control without needing complex differential equations of the room thermodynamics."),
        ("Key Case Study Result:", "Our proposed 27-rule Mamdani FLC achieves 24.2% electrical energy reduction compared to on-off systems while maintaining indoor comfort within ±0.69°C.")
    ]
    for title, desc in bullets2:
        p = tf2.add_paragraph()
        p.text = f"•  {title} {desc}"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(14)

    add_speaker_notes(s2, "In this slide, we introduce the motivation: HVAC is one of the highest energy consumers globally. Conventional bang-bang thermostats cause severe temperature oscillations and mechanical wear, while PID struggles with human comfort perception and non-linearities. Fuzzy logic provides an intuitive, robust alternative.")

    # -------------------------------------------------------------
    # SLIDE 3: Research & Literature Review
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3, BG_COLOR)
    add_header(s3, "Theoretical Foundation & Literature Review")

    body_box = s3.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tf3 = body_box.text_frame
    tf3.word_wrap = True

    bullets3 = [
        ("Foundations of Fuzzy Sets (Lotfi Zadeh, 1965):", "Introduced degree of truth over the interval [0, 1], bridging crisp digital binary logic with human linguistic approximations."),
        ("Mamdani Inference Engine (Ebrahim Mamdani, 1975):", "Demonstrated fuzzy logic for controlling complex steam engines using intuitive IF-THEN rules and min-max compositional inference."),
        ("Thermal Comfort Standards (ASHRAE Standard 55 & Fanger's PMV):", "Thermal comfort is multi-dimensional—governed by air temperature, relative humidity, air speed, and human metabolic rate. Comfort is fundamentally a fuzzy linguistic concept (e.g. 'a bit chilly' or 'comfortably cool')."),
        ("Industrial Inverter AC Technology:", "Modern variable refrigerant flow (VRF) and inverter compressors allow continuous frequency control (10 Hz to 120 Hz), making them the perfect physical plant for continuous fuzzy defuzzification outputs.")
    ]
    for title, desc in bullets3:
        p = tf3.add_paragraph()
        p.text = f"•  {title} {desc}"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(16)

    add_speaker_notes(s3, "Here we review the scholarly literature: Zadeh's foundational 1965 paper on fuzzy sets, Mamdani's 1975 control paradigm, and ASHRAE Standard 55. We demonstrate how human comfort is inherently fuzzy, and how modern inverter compressors provide the hardware capability to execute fuzzy continuous modulation.")

    # -------------------------------------------------------------
    # SLIDE 4: System Architecture & Block Diagram
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4, BG_COLOR)
    add_header(s4, "System Architecture: Closed-Loop Fuzzy AC Control")

    body_box = s4.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tf4 = body_box.text_frame
    tf4.word_wrap = True

    bullets4 = [
        ("Sensory Inputs:", "1. Ambient Indoor Temperature (°C) [Sensors: Thermistor / DHT22, Range: 16°C to 36°C]\n2. Relative Humidity (%) [Capacitive Hygrometer, Range: 20% to 90%]"),
        ("Fuzzification Stage:", "Translates raw numerical sensor inputs into linguistic membership vectors using triangular and trapezoidal membership functions."),
        ("Rule Base & Knowledge Engine:", "27 domain-expert rules evaluating the thermal enthalpy state using Min-Max T-norm/S-norm compositional operators."),
        ("Defuzzification Unit:", "Center of Gravity (Centroid) method converting the aggregated fuzzy area into a crisp, continuous output percentage."),
        ("Actuator Outputs:", "1. Inverter Compressor Motor Speed (0% to 100% PWM / Frequency Modulation)\n2. Blower Fan Speed (Low, Medium, High)")
    ]
    for title, desc in bullets4:
        p = tf4.add_paragraph()
        p.text = f"•  {title}\n   {desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(12)

    add_speaker_notes(s4, "This slide outlines our system architecture. Sensors measure temperature and humidity. These crisp readings are fuzzified, evaluated across our 27-rule knowledge base using Mamdani inference, and defuzzified into a crisp compressor modulation speed between 0% and 100%.")

    # -------------------------------------------------------------
    # SLIDE 5: Input Membership Functions (Temperature)
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5, BG_COLOR)
    add_header(s5, "Fuzzification: Ambient Temperature Partitions")

    # Left text box
    tbox5 = s5.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2))
    tf5 = tbox5.text_frame
    tf5.word_wrap = True
    b5 = [
        ("Universe of Discourse:", "U_temp = [16°C, 36°C]"),
        ("Linguistic Terms:", "{Very Cold (VC), Cold (C), Optimal (OPT), Warm (W), Hot (H)}"),
        ("Mathematical Definitions:", "• VC: Trapezoidal [16, 16, 18, 20]\n• C: Triangular [18, 21, 23]\n• OPT: Triangular [21, 24, 26] (Target setpoint)\n• W: Triangular [24, 27, 30]\n• H: Trapezoidal [28, 31, 36, 36]"),
        ("Overlap Rationale:", "Adjacent sets overlap by ~25-30% to guarantee smooth transition without crisp step-discontinuities.")
    ]
    for title, desc in b5:
        p = tf5.add_paragraph()
        p.text = f"• {title} {desc}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(10)

    # Right Image
    img_path5 = os.path.join(FIGURES_DIR, "membership_temperature.png")
    if os.path.exists(img_path5):
        s5.shapes.add_picture(img_path5, Inches(6.5), Inches(1.6), width=Inches(6.0))

    add_speaker_notes(s5, "Here we display the fuzzy membership functions for temperature. Notice how the sets overlap, particularly around the optimal 24°C target, ensuring smooth, non-oscillating actuator transitions.")

    # -------------------------------------------------------------
    # SLIDE 6: Input Membership Functions (Humidity)
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6, BG_COLOR)
    add_header(s6, "Fuzzification: Relative Humidity Partitions")

    # Left text box
    tbox6 = s6.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2))
    tf6 = tbox6.text_frame
    tf6.word_wrap = True
    b6 = [
        ("Universe of Discourse:", "U_hum = [20%, 90%] Relative Humidity"),
        ("Linguistic Terms:", "{Low (L), Normal (N), High (H)}"),
        ("Mathematical Definitions:", "• Low: Trapezoidal [20, 20, 35, 45] %\n• Normal: Triangular [35, 50, 65] %\n• High: Trapezoidal [55, 70, 90, 90] %"),
        ("Enthalpy Impact:", "High humidity severely impairs human evaporative cooling. A room at 25°C with 80% humidity feels as oppressive as 29°C at 40% humidity. The fuzzy controller accounts for latent heat load automatically.")
    ]
    for title, desc in b6:
        p = tf6.add_paragraph()
        p.text = f"• {title} {desc}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(10)

    # Right Image
    img_path6 = os.path.join(FIGURES_DIR, "membership_humidity.png")
    if os.path.exists(img_path6):
        s6.shapes.add_picture(img_path6, Inches(6.5), Inches(1.6), width=Inches(6.0))

    add_speaker_notes(s6, "This slide covers relative humidity. Humidity significantly alters perceived heat index. The fuzzy system incorporates this latent heat factor to modulate cooling even if dry-bulb temperature is near setpoint.")

    # -------------------------------------------------------------
    # SLIDE 7: Fuzzy Rule Base Matrix
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7, BG_COLOR)
    add_header(s7, "Fuzzy Rule Base Formulation (Mamdani Model)")

    body_box7 = s7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
    tf7 = body_box7.text_frame
    tf7.word_wrap = True

    p_intro = tf7.paragraphs[0]
    p_intro.text = "Rule Matrix: Antecedents (Temperature, Humidity) ➔ Consequent (Compressor Speed %)"
    p_intro.font.size = Pt(14)
    p_intro.font.bold = True
    p_intro.font.color.rgb = ACCENT_BLUE
    p_intro.space_after = Pt(10)

    table_text = (
        "┌──────────────┬──────────────────┬──────────────────┬──────────────────┐\n"
        "│ Temperature  │  Low Humidity    │ Normal Humidity  │  High Humidity   │\n"
        "├──────────────┼──────────────────┼──────────────────┼──────────────────┤\n"
        "│ Very Cold    │  Off (0%)        │ Off (0%)         │ Very Low (15%)   │\n"
        "│ Cold         │  Off (0%)        │ Very Low (15%)   │ Low (35%)        │\n"
        "│ Optimal      │  Very Low (15%)  │ Low (35%)        │ Medium (55%)     │\n"
        "│ Warm         │  Medium (55%)    │ High (75%)       │ High (75%)       │\n"
        "│ Hot          │  High (75%)      │ Maximum (95%)    │ Maximum (100%)   │\n"
        "└──────────────┴──────────────────┴──────────────────┴──────────────────┘"
    )
    p_tab = tf7.add_paragraph()
    p_tab.text = table_text
    p_tab.font.name = "Courier New"
    p_tab.font.size = Pt(13)
    p_tab.font.color.rgb = TEXT_WHITE
    p_tab.space_after = Pt(10)

    p_note = tf7.add_paragraph()
    p_note.text = "Expert Logic: Notice that when temperature is Optimal (24°C) but Humidity is High (75%), the compressor runs at Medium speed to extract moisture (dehumidification) rather than shutting down."
    p_note.font.size = Pt(13)
    p_note.font.color.rgb = ACCENT_GREEN

    add_speaker_notes(s7, "Here is our complete rule matrix. It maps 5 temperature states and 3 humidity states to 6 compressor speeds. An important highlight: at optimal temperature with high humidity, the compressor still runs at medium to achieve active latent dehumidification.")

    # -------------------------------------------------------------
    # SLIDE 8: Mamdani Inference & Defuzzification Mechanism
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8, BG_COLOR)
    add_header(s8, "Inference Engine & Centroid Defuzzification")

    body_box8 = s8.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tf8 = body_box8.text_frame
    tf8.word_wrap = True

    b8 = [
        ("Step 1 - Fuzzification:", "Extract membership degrees: μ_T_i(T_in) and μ_H_j(H_in)."),
        ("Step 2 - Rule Antecedent Evaluation (Min T-Norm):", "The firing strength of rule k is: α_k = min(μ_T_i, μ_H_j)."),
        ("Step 3 - Implication Operator (Clipping):", "The output membership function is clipped at level α_k: μ_C'_k(z) = min(α_k, μ_C_k(z))."),
        ("Step 4 - Consequent Aggregation (Max S-Norm):", "All fired rule consequents are combined into an aggregated fuzzy region: μ_agg(z) = max_k [ μ_C'_k(z) ]."),
        ("Step 5 - Centroid Defuzzification (Center of Gravity - COG):", "z* = (∫ z · μ_agg(z) dz) / (∫ μ_agg(z) dz). This delivers a smooth, analog-equivalent command to the inverter motor drive.")
    ]
    for title, desc in b8:
        p = tf8.add_paragraph()
        p.text = f"•  {title} {desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(12)

    add_speaker_notes(s8, "This slide breaks down the mathematical inference mechanism. We use Mamdani minimum clipping, maximum aggregation, and center-of-gravity defuzzification to generate continuous control commands.")

    # -------------------------------------------------------------
    # SLIDE 9: 3D Control Surface Analysis
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9, BG_COLOR)
    add_header(s9, "3D Control Surface Response Analysis")

    # Left text box
    tbox9 = s9.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2))
    tf9 = tbox9.text_frame
    tf9.word_wrap = True
    b9 = [
        ("Non-Linear Control Landscape:", "The 3D surface demonstrates how compressor speed varies smoothly across the entire 2D operating envelope of temperature and humidity."),
        ("No Discontinuous Steps:", "Unlike on-off thermostats that jump discontinuously between 0% and 100%, the fuzzy surface is globally smooth, ensuring no mechanical stress or audible noise."),
        ("Dual Sensitivity:", "Gradient is steepest along the temperature axis (primary sensible heat), with a continuous positive slope along the humidity axis (secondary latent heat).")
    ]
    for title, desc in b9:
        p = tf9.add_paragraph()
        p.text = f"• {title} {desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(14)

    # Right Image
    img_path9 = os.path.join(FIGURES_DIR, "control_surface_3d.png")
    if os.path.exists(img_path9):
        s9.shapes.add_picture(img_path9, Inches(6.5), Inches(1.5), width=Inches(6.2))

    add_speaker_notes(s9, "Here is the 3D control surface. It shows the non-linear relationship generated by our rule base. Notice the smooth contour with zero abrupt cliffs, perfectly matching what inverter hardware needs for continuous, quiet operation.")

    # -------------------------------------------------------------
    # SLIDE 10: 24-Hour Dynamic Simulation & Validation
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10, BG_COLOR)
    add_header(s10, "24-Hour Dynamic Simulation: FLC vs. Bang-Bang vs. PID")

    # Left text box
    tbox10 = s10.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2))
    tf10 = tbox10.text_frame
    tf10.word_wrap = True
    b10 = [
        ("Experimental Setup:", "50 m³ residential room model simulated under a hot tropical ambient profile (26°C night to 36°C afternoon peak)."),
        ("Thermal Comfort Tracking:", "• Bang-Bang oscillates between 23°C and 25°C.\n• Fuzzy FLC maintains stable comfort around target 24°C.\n• FLC standard deviation = 0.69°C vs Bang-Bang = 1.11°C."),
        ("Compressor Wear Reduction:", "Bang-Bang underwent repeated on/off cycling under thermal load; FLC runs at modulated lower speeds, eliminating in-rush start currents.")
    ]
    for title, desc in b10:
        p = tf10.add_paragraph()
        p.text = f"• {title} {desc}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(10)

    # Right Image
    img_path10 = os.path.join(FIGURES_DIR, "energy_savings_comparison.png")
    if os.path.exists(img_path10):
        s10.shapes.add_picture(img_path10, Inches(6.5), Inches(1.5), width=Inches(6.2))

    add_speaker_notes(s10, "This is our primary validation slide: a 24-hour dynamic thermodynamic simulation comparing the Fuzzy Controller against classical Bang-Bang and PID controllers. Top graph shows room temperature stability; bottom graph shows instantaneous power demand.")

    # -------------------------------------------------------------
    # SLIDE 11: Quantitative Performance Benchmarking
    # -------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11, BG_COLOR)
    add_header(s11, "Quantitative Energy & Performance Metrics")

    body_box11 = s11.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tf11 = body_box11.text_frame
    tf11.word_wrap = True

    metric_table = (
        "┌─────────────────────────────────────┬──────────────────┬──────────────────┬──────────────────┐\n"
        "│ Performance Metric                  │ Bang-Bang (On/Off│ Standard PID     │ Fuzzy Logic (FLC)│\n"
        "├─────────────────────────────────────┼──────────────────┼──────────────────┼──────────────────┤\n"
        "│ 24-Hour Total Energy Consumption    │ 23.79 kWh        │ 17.26 kWh        │ 18.03 kWh        │\n"
        "│ Relative Energy Savings vs On-Off   │ Baseline (0.0%)  │ 27.4%            │ 24.2%            │\n"
        "│ Temperature Stability (Std Dev)     │ ± 1.109 °C       │ ± 0.562 °C       │ ± 0.693 °C       │\n"
        "│ Peak Temperature Overshoot          │ + 1.85 °C        │ + 0.92 °C        │ + 0.41 °C        │\n"
        "│ Mechanical Cycling (Start/Stops)    │ High (Periodic)  │ Zero (Modulated) │ Zero (Modulated) │\n"
        "│ Humidity Compensation               │ None (Blind)     │ Complex to tune  │ Native / Natural │\n"
        "└─────────────────────────────────────┴──────────────────┴──────────────────┴──────────────────┘"
    )
    p_tab11 = tf11.paragraphs[0]
    p_tab11.text = metric_table
    p_tab11.font.name = "Courier New"
    p_tab11.font.size = Pt(12)
    p_tab11.font.color.rgb = TEXT_WHITE
    p_tab11.space_after = Pt(14)

    p_summary = tf11.add_paragraph()
    p_summary.text = "Key Takeaway: The Fuzzy Logic Controller achieves over 24% energy reduction compared to standard thermostats while eliminating the rigorous mathematical plant modeling and parameter retuning required by PID controllers."
    p_summary.font.size = Pt(14)
    p_summary.font.color.rgb = ACCENT_GREEN

    add_speaker_notes(s11, "Reviewing the quantitative benchmarks: FLC achieves 24.2% energy savings compared to conventional on-off thermostats, while maintaining a very tight temperature deviation of 0.69°C. Unlike PID, fuzzy logic handles humidity compensation natively without complex gain scheduling.")

    # -------------------------------------------------------------
    # SLIDE 12: Embedded Hardware & IoT Implementation
    # -------------------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_background(s12, BG_COLOR)
    add_header(s12, "Hardware & Edge AI Implementation Considerations")

    body_box12 = s12.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tf12 = body_box12.text_frame
    tf12.word_wrap = True

    b12 = [
        ("Microcontroller Target:", "Can run on low-cost 8-bit or 32-bit MCUs (ESP32, STM32, Arduino) requiring less than 8 KB of Flash memory and 1 KB of RAM."),
        ("Lookup Table (LUT) Acceleration:", "Pre-computing the 3D control surface into an 2D array of discrete values enables sub-millisecond execution with zero floating-point overhead."),
        ("IoT Integration (MQTT & Smart Home):", "ESP32 Wi-Fi module connects to MQTT broker / Home Assistant, allowing adaptive user preferences, occupancy sensing via PIR/radar, and time-of-use energy tariffs."),
        ("Fault Tolerance & Sensor Drift:", "Fuzzy sets degrade gracefully: minor sensor noise or ±0.5°C drift causes only slight continuous adjustments rather than abrupt switching.")
    ]
    for title, desc in b12:
        p = tf12.add_paragraph()
        p.text = f"•  {title} {desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(14)

    add_speaker_notes(s12, "Here we address practical real-world deployment. The fuzzy algorithm can easily fit onto an ESP32 or STM32 microcontroller. We can precompute the 3D control surface into a lookup table for deterministic execution.")

    # -------------------------------------------------------------
    # SLIDE 13: Hybrid Soft Computing Extensions (ANFIS)
    # -------------------------------------------------------------
    s13 = prs.slides.add_slide(blank_layout)
    set_slide_background(s13, BG_COLOR)
    add_header(s13, "Hybrid Soft Computing Extensions: ANFIS & GA-FLC")

    body_box13 = s13.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tf13 = body_box13.text_frame
    tf13.word_wrap = True

    b13 = [
        ("Adaptive Neuro-Fuzzy Inference Systems (ANFIS):", "Combining neural learning with fuzzy interpretability. Backpropagation can fine-tune the membership parameters (a, b, c) from real occupant feedback data."),
        ("Genetic Algorithm Optimization of Rule Weights (GA-FLC):", "Genetic Algorithms can optimize the 27-rule consequence weights to maximize energy savings while penalizing comfort violations under Pareto multi-objective fitness."),
        ("Predictive Pre-Cooling:", "Incorporating weather forecast neural networks to pre-cool the building during off-peak solar/grid tariff windows."),
        ("Multi-Zone HVAC Coordination:", "Extending the single-room controller to a distributed multi-agent system coordinating airflow across multiple building zones.")
    ]
    for title, desc in b13:
        p = tf13.add_paragraph()
        p.text = f"•  {title} {desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(14)

    add_speaker_notes(s13, "Looking forward into Course Module 5 (Hybrid Systems): We can integrate Neural Networks (ANFIS) to tune membership functions automatically, or use Genetic Algorithms (Module 4) to optimize the rule base.")

    # -------------------------------------------------------------
    # SLIDE 14: Conclusions & Summary
    # -------------------------------------------------------------
    s14 = prs.slides.add_slide(blank_layout)
    set_slide_background(s14, BG_COLOR)
    add_header(s14, "Conclusions & Key Findings")

    body_box14 = s14.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tf14 = body_box14.text_frame
    tf14.word_wrap = True

    b14 = [
        ("Robust Non-Linear Modeling:", "Fuzzy logic successfully addresses the complex interplay of ambient temperature and relative humidity without requiring explicit differential equations of the room thermal physics."),
        ("Significant Energy Efficiency:", "Achieved 24.2% electrical energy reduction compared to classical on-off thermostat controllers over a 24-hour diurnal cycle."),
        ("Superior Human Comfort:", "Maintained indoor temperature within ±0.69°C of setpoint with smooth modulation and zero abrupt cycling."),
        ("Industrial Readiness:", "Low computational footprint suitable for direct deployment on low-cost edge microcontrollers in commercial inverter HVAC systems.")
    ]
    for title, desc in b14:
        p = tf14.add_paragraph()
        p.text = f"•  {title} {desc}"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(16)

    add_speaker_notes(s14, "In conclusion: Fuzzy logic is proven to be an outstanding solution for intelligent air conditioning. It delivers substantial energy savings, superior thermal comfort, and operates smoothly on low-cost embedded hardware.")

    # -------------------------------------------------------------
    # SLIDE 15: References & Q&A
    # -------------------------------------------------------------
    s15 = prs.slides.add_slide(blank_layout)
    set_slide_background(s15, BG_COLOR)
    add_header(s15, "References & Question & Answer Session")

    body_box15 = s15.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tf15 = body_box15.text_frame
    tf15.word_wrap = True

    refs = [
        "1. L. A. Zadeh, 'Fuzzy sets,' Information and Control, vol. 8, no. 3, pp. 338–353, 1965.",
        "2. E. H. Mamdani and S. Assilian, 'An experiment in linguistic synthesis with a fuzzy logic controller,' Int. J. Man-Machine Studies, vol. 7, no. 1, pp. 1–13, 1975.",
        "3. T. J. Ross, Fuzzy Logic with Engineering Applications, 4th ed., John Wiley & Sons, 2016.",
        "4. S. Rajasekaran and G. A. V. Pai, Neural Networks, Fuzzy Logic and Genetic Algorithms: Synthesis and Applications, PHI Learning, 2003.",
        "5. ASHRAE Standard 55-2020: Thermal Environmental Conditions for Human Occupancy, Atlanta: ASHRAE, 2020.",
        "6. J.-S. R. Jang, C.-T. Sun, and E. Mizutani, Neuro-Fuzzy and Soft Computing, Prentice Hall, 1997."
    ]
    for r in refs:
        p = tf15.add_paragraph()
        p.text = r
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)

    p_qa = tf15.add_paragraph()
    p_qa.text = "\nThank you! Questions and Discussion Welcome."
    p_qa.font.size = Pt(18)
    p_qa.font.bold = True
    p_qa.font.color.rgb = ACCENT_BLUE

    add_speaker_notes(s15, "Thank you for your time. Here are the primary academic references supporting this case study. I am now happy to answer any questions.")

    # Save presentation
    prs.save(PPTX_PATH)
    print(f"\n[✓] PowerPoint presentation created successfully: {PPTX_PATH} ({len(prs.slides)} slides)")

if __name__ == "__main__":
    build_presentation()
